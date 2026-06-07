from fastapi import FastAPI, HTTPException, BackgroundTasks, UploadFile, File, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import List, Optional, Union
from datetime import datetime, timedelta
import uuid
import time
import asyncio
import sqlite3
import pickle
import json
import re
import hashlib
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.embeddings import Embeddings
from rank_bm25 import BM25Okapi
import os
import shutil
import threading
import uvicorn
import requests
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), ".env"))

app = FastAPI(title="知识库API（检索优化版）")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================
# 配置常量
# ============================================================
MAX_HISTORY_EXCHANGES = 10
SESSION_MAX_AGE_DAYS = 7
SUMMARY_TRIGGER_COUNT = 10

# 检索优化配置
VECTOR_WEIGHT = 0.7       # 混合检索中向量检索权重
BM25_WEIGHT = 0.3          # 混合检索中 BM25 权重
HYBRID_CANDIDATE_K = 15    # 混合检索候选数量
DEFAULT_TOP_K = 3          # 最终返回文档数

# ============================================================
# Pydantic 模型
# ============================================================

def _normalize_pre_retrieval(v):
    """将 pre_retrieval 统一转为有序列表，保证执行顺序: rewrite -> hyde -> mqe"""
    ORDER = ["rewrite", "hyde", "mqe"]
    if v is None or v == "none":
        return []
    if isinstance(v, str):
        return [] if v == "none" else [v]
    if isinstance(v, list):
        cleaned = [x for x in v if x and x != "none"]
        return [s for s in ORDER if s in cleaned]
    return []


class Query(BaseModel):
    question: str
    session_id: Optional[str] = None
    kb_id: Optional[str] = None
    retrieval_strategy: Optional[str] = "default"   # "default" | "hybrid"
    pre_retrieval: Optional[Union[str, List[str]]] = "none"
    post_retrieval: Optional[str] = "none"           # "none" | "rerank"
    top_k: Optional[int] = None

    def __init__(self, **data):
        if "pre_retrieval" in data:
            data["pre_retrieval"] = _normalize_pre_retrieval(data["pre_retrieval"])
        if "top_k" in data and data["top_k"] is not None:
            data["top_k"] = max(1, min(data["top_k"], 20))
        super().__init__(**data)

class Response(BaseModel):
    answer: str
    sources: list = []
    session_id: Optional[str] = None

class OpenAIModel(BaseModel):
    id: str
    object: str = "model"
    created: int = Field(default_factory=lambda: int(time.time()))
    owned_by: str = "knowledge-base"

class OpenAIModelsResponse(BaseModel):
    object: str = "list"
    data: List[OpenAIModel]

class OpenAIMessage(BaseModel):
    role: str
    content: str

class OpenAIChatRequest(BaseModel):
    model: str = "knowledge-base"
    messages: List[OpenAIMessage]
    temperature: Optional[float] = 0.7
    max_tokens: Optional[int] = None
    stream: Optional[bool] = False
    session_id: Optional[str] = None
    kb_id: Optional[str] = None
    retrieval_strategy: Optional[str] = "default"
    pre_retrieval: Optional[Union[str, List[str]]] = "none"
    post_retrieval: Optional[str] = "none"
    top_k: Optional[int] = None

    def __init__(self, **data):
        if "pre_retrieval" in data:
            data["pre_retrieval"] = _normalize_pre_retrieval(data["pre_retrieval"])
        if "top_k" in data and data["top_k"] is not None:
            data["top_k"] = max(1, min(data["top_k"], 20))
        super().__init__(**data)

class OpenAIChatChoice(BaseModel):
    index: int = 0
    message: OpenAIMessage
    finish_reason: str = "stop"

class OpenAIUsage(BaseModel):
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0

class OpenAIChatResponse(BaseModel):
    id: str = Field(default_factory=lambda: f"chatcmpl-{uuid.uuid4().hex[:24]}")
    object: str = "chat.completion"
    created: int = Field(default_factory=lambda: int(time.time()))
    model: str = "knowledge-base"
    choices: List[OpenAIChatChoice]
    usage: OpenAIUsage = OpenAIUsage()

# 会话相关模型
class SessionCreate(BaseModel):
    title: Optional[str] = "新对话"
    kb_id: Optional[str] = "default"

class SessionInfo(BaseModel):
    session_id: str
    kb_id: str = "default"
    title: str
    created_at: str
    updated_at: str
    summary: Optional[str] = None
    message_count: int = 0

# 知识库相关模型
class KBCreate(BaseModel):
    name: str
    description: Optional[str] = ""

class KBInfo(BaseModel):
    id: str
    name: str
    created_at: str
    updated_at: str
    description: str = ""
    file_count: int = 0

class SessionListResponse(BaseModel):
    sessions: List[SessionInfo]

class MessageInfo(BaseModel):
    role: str
    content: str
    created_at: str

class SessionHistoryResponse(BaseModel):
    session_id: str
    messages: List[MessageInfo]

# ============================================================
# Prompt 模板
# ============================================================

prompt_template = """你是一个专业的 AI 助手。你的核心职责是基于提供的参考资料回答用户问题。

### 回答原则：
1. 严禁编造知识库中的信息。
2. 无论问题是否与参考资料相关，你都必须给出有益的回复，但同时必须严格区分并标注信息的来源。

### 来源标注规范（必须严格执行）：
- **来自知识库**：陈述参考资料中的内容时，必须在对应句子或段落末尾标注，如 [来源:xx文件xx节]。
- **超出知识库**：如果问题无法从参考资料中找到答案，你必须先明确声明"该问题未在知识库中找到相关资料"，然后可以调用你的通用知识进行补充解答，并在补充内容后标注 [来源:通用知识]。
- **混合情况**：如果回答中既有参考资料的内容，又有你补充的通用知识，必须分别标注，绝不能混淆。

### 回答格式要求：
1. 语言简洁，逻辑清晰，使用列表或分段提升可读性。

---
参考资料：
{context}
---
用户问题：
{question}
---
### 回答"""

PROMPT = PromptTemplate(
    template=prompt_template,
    input_variables=["context", "question"]
)

prompt_template_with_history = """你是一个专业的 AI 助手。你的核心职责是基于提供的参考资料回答用户问题，并结合对话上下文理解意图。

### 回答原则：
1. 严禁编造知识库中的信息。
2. 无论问题是否与参考资料相关，你都必须给出有益的回复，但同时必须严格区分并标注信息的来源。

### 来源标注规范（必须严格执行）：
- **来自知识库**：陈述参考资料中的内容时，必须在对应句子或段落末尾标注，如 [来源:xx文件xx节]。
- **超出知识库**：如果问题无法从参考资料中找到答案，你必须先明确声明"该问题未在知识库中找到相关资料"，然后可以调用你的通用知识进行补充解答，并在补充内容后标注 [来源:通用知识]。
- **混合情况**：如果回答中既有参考资料的内容，又有你补充的通用知识，必须分别标注，绝不能混淆。

### 回答格式要求：
1. 语言简洁，逻辑清晰，使用列表或分段提升可读性。
2. 结合对话上下文，准确理解代词和省略的指代对象（如"它"、"这个方法"、"请详细解释"等）。

---
参考资料：
{context}
---
历史对话：
{history}
---
用户问题：
{question}
---
### 回答："""

PROMPT_WITH_HISTORY = PromptTemplate(
    template=prompt_template_with_history,
    input_variables=["context", "history", "question"]
)

SUMMARY_PROMPT_TEMPLATE = "请用2-3句话总结以下对话的主要内容：\n\n{conversation}\n\n总结："

QUERY_REWRITE_PROMPT = """你是一个检索优化助手。请将以下用户问题改写为更适合向量检索的规范化查询词。
要求：
1. 去除口语化表达和代词（如"它"、"这个"、"那"）
2. 补充缺失的主语（结合对话历史推断代词指代）
3. 保留核心专业术语
4. 只输出改写后的查询词，不要解释

对话历史：
{history}

用户问题：{question}

改写后的查询词："""

HYDE_PROMPT = """请根据以下问题，写一段简短的假设性答案（约100字）。
这段答案不需要准确，只需包含可能出现在文档中的专业术语和表述方式。

问题：{question}

假设性答案："""

MQE_PROMPT = """你是一个检索优化助手。请根据以下用户问题，生成 {n} 个不同的检索查询变体。
每个变体应从不同角度或使用不同表述来覆盖原始问题的语义。
要求：
1. 每个查询变体单独一行
2. 只输出查询词，不要编号、不要解释
3. 变体之间应有差异（同义替换、角度切换、简化/详细化等）

用户问题：{question}

查询变体："""

# ============================================================
# 全局变量
# ============================================================
embeddings = None
llm = None
reranker = None

# 硅基流动共享配置（所有 SF 模型共用 base_url 和 api_key）
SILICONFLOW_BASE_URL = "https://api.siliconflow.cn/v1"
siliconflow_api_key = os.getenv("SILICONFLOW_API_KEY") or os.getenv("EMBEDDING_API_KEY") or os.getenv("LLM_API_KEY")

# 各模型配置（默认全部使用硅基流动，但尊重 .env 中的自定义值）
embedding_base_url = os.getenv("EMBEDDING_BASE_URL") or SILICONFLOW_BASE_URL
embedding_model_name = os.getenv("EMBEDDING_MODEL", "BAAI/bge-m3")

reranker_base_url = os.getenv("RERANKER_BASE_URL") or SILICONFLOW_BASE_URL
RERANKER_MODEL = os.getenv("RERANKER_MODEL", "BAAI/bge-reranker-v2-m3")

llm_base_url = os.getenv("LLM_BASE_URL") or SILICONFLOW_BASE_URL
llm_model = os.getenv("LLM_MODEL", "Qwen/Qwen3-8B")

# API key: 硅基流动 URL 使用共享 key，自定义 URL 使用各自的 key
def _resolve_startup_key(base_url, env_key_name):
    if "siliconflow" in (base_url or ""):
        return siliconflow_api_key
    return os.getenv(env_key_name)

embedding_api_key = _resolve_startup_key(embedding_base_url, "EMBEDDING_API_KEY")
reranker_api_key = _resolve_startup_key(reranker_base_url, "RERANKER_API_KEY")
llm_api_key = _resolve_startup_key(llm_base_url, "LLM_API_KEY")

# 数据库路径
base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DB_PATH = os.path.join(base_dir, "backend", "sessions.db")

# ============================================================
# OpenAICompatibleEmbeddings 类
# ============================================================

class OpenAICompatibleEmbeddings(Embeddings):
    """OpenAI 兼容嵌入模型（默认硅基流动 BAAI/bge-m3）"""
    def __init__(self, api_key=None, base_url=None, model=None):
        self.api_key = api_key or embedding_api_key
        self.base_url = base_url or embedding_base_url
        self.model = model or embedding_model_name

    def _get_embeddings(self, texts_list):
        url = f"{self.base_url}/embeddings"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}"
        }
        data = {
            "model": self.model,
            "input": texts_list
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            result = response.json()
            embeddings = [item['embedding'] for item in result['data']]
            return embeddings
        except Exception as e:
            print(f"调用嵌入模型API失败: {e}")
            print(f"响应内容: {response.text if 'response' in locals() else '无响应'}")
            raise RuntimeError(f"嵌入模型API调用失败: {e}") from e

    def embed_documents(self, texts):
        return self._get_embeddings(texts)

    def embed_query(self, text):
        result = self._get_embeddings([text])
        return result[0]

# ============================================================
# OpenAICompatibleReranker 类（后检索优化）
# ============================================================

class OpenAICompatibleReranker:
    """OpenAI 兼容重排序模型（默认硅基流动 BAAI/bge-reranker-v2-m3）"""
    def __init__(self, api_key=None, base_url=None, model=None, top_n=DEFAULT_TOP_K):
        self.api_key = api_key or reranker_api_key
        self.base_url = base_url or reranker_base_url
        self.model = model or RERANKER_MODEL
        self.top_n = top_n

    def rerank(self, query: str, documents: list, top_n: int = None) -> list:
        if not documents:
            return []
        top_n = top_n or self.top_n

        doc_texts = [doc.page_content if hasattr(doc, 'page_content') else str(doc) for doc in documents]

        url = f"{self.base_url}/rerank"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}"
        }
        data = {
            "model": self.model,
            "query": query,
            "documents": doc_texts,
            "top_n": min(top_n, len(doc_texts)),
            "return_documents": False
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=60)
            response.raise_for_status()
            result = response.json()
            reranked = []
            for item in result['results']:
                idx = item['index']
                score = item['relevance_score']
                doc = documents[idx]
                # Attach rerank score to metadata
                if hasattr(doc, 'metadata'):
                    doc.metadata['rerank_score'] = round(score, 4)
                reranked.append(doc)
            return reranked
        except Exception as e:
            print(f"重排序API调用失败: {e}")
            print(f"响应内容: {response.text if 'response' in locals() else '无响应'}")
            return documents[:top_n]

# ============================================================
# 中文分词工具（BM25 用）
# ============================================================

def tokenize(text: str) -> list:
    """中英文混合分词：中文按字符 bigram，英文按空格"""
    text = text.lower().strip()
    tokens = []
    # 英文单词
    en_tokens = re.findall(r'[a-zA-Z]+', text)
    tokens.extend(en_tokens)
    # 中文字符 bigram
    cn_chars = re.findall(r'[一-鿿]', text)
    for i in range(len(cn_chars) - 1):
        tokens.append(cn_chars[i] + cn_chars[i + 1])
    # 单个中文字符也保留
    tokens.extend(cn_chars)
    return tokens

# ============================================================
# 多知识库管理器
# ============================================================

class KBManager:
    """管理多个知识库，支持懒加载和 LRU 缓存淘汰"""

    def __init__(self, max_loaded=2):
        self._loaded = {}          # kb_id -> {vectorstore, bm25_index, bm25_docs, retriever}
        self._max_loaded = max_loaded
        self._access_order = []    # LRU 追踪
        self._lock = threading.Lock()

    @staticmethod
    def _to_slug(kb_id: str) -> str:
        """将 kb_id 转为纯 ASCII 目录名，避免 FAISS C++ 后端不支持 Unicode 路径"""
        ascii_part = re.sub(r'[^a-zA-Z0-9]', '', kb_id)
        if ascii_part:
            return ascii_part.lower()
        return f"kb{hashlib.md5(kb_id.encode('utf-8')).hexdigest()[:8]}"

    def get_kb_path(self, kb_id: str) -> str:
        slug = self._to_slug(kb_id)
        slug_path = os.path.join(base_dir, "knowledge_bases", slug)
        raw_path = os.path.join(base_dir, "knowledge_bases", kb_id)
        # 非 ASCII 路径 → 优先迁移（保证后续所有操作用同一路径）
        if kb_id != slug and os.path.isdir(raw_path):
            try:
                os.rename(raw_path, slug_path)
                print(f"知识库目录已迁移: '{kb_id}' -> '{slug}'")
            except Exception as e:
                print(f"目录迁移失败: {e}，使用原路径")
                return raw_path
        return slug_path

    def get_materials_path(self, kb_id: str) -> str:
        return os.path.join(self.get_kb_path(kb_id), "materials")

    def registry_path(self) -> str:
        return os.path.join(base_dir, "knowledge_bases", "registry.json")

    def load_registry(self) -> dict:
        path = self.registry_path()
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        return {}

    def save_registry(self, registry: dict):
        path = self.registry_path()
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            json.dump(registry, f, ensure_ascii=False, indent=2)

    def list_kbs(self) -> list:
        registry = self.load_registry()
        result = []
        for kb_id, info in registry.items():
            kb_path = self.get_kb_path(kb_id)
            materials_path = self.get_materials_path(kb_id)
            file_count = 0
            if os.path.isdir(materials_path):
                file_count = len([f for f in os.listdir(materials_path)
                                  if os.path.isfile(os.path.join(materials_path, f)) and not f.startswith(".")])
            result.append({**info, "file_count": file_count})
        return result

    def create_kb(self, name: str, description: str = "") -> dict:
        slug = self._to_slug(name)
        registry = self.load_registry()
        if slug in registry:
            slug = f"{slug}{int(time.time()) % 10000}"
        now = datetime.now().isoformat()
        info = {"id": slug, "name": name, "created_at": now, "updated_at": now, "description": description}
        registry[slug] = info
        self.save_registry(registry)
        kb_path = self.get_kb_path(slug)
        os.makedirs(os.path.join(kb_path, "materials"), exist_ok=True)
        return info

    def delete_kb(self, kb_id: str) -> bool:
        if kb_id == "default":
            return False
        registry = self.load_registry()
        if kb_id not in registry:
            return False
        del registry[kb_id]
        self.save_registry(registry)
        self.invalidate(kb_id)
        kb_path = self.get_kb_path(kb_id)
        if os.path.isdir(kb_path):
            shutil.rmtree(kb_path)
        return True

    def update_kb(self, kb_id: str, name: str = None, description: str = None) -> bool:
        registry = self.load_registry()
        if kb_id not in registry:
            return False
        if name is not None:
            registry[kb_id]["name"] = name
        if description is not None:
            registry[kb_id]["description"] = description
        registry[kb_id]["updated_at"] = datetime.now().isoformat()
        self.save_registry(registry)
        return True

    def get(self, kb_id: str) -> dict:
        with self._lock:
            if kb_id in self._loaded:
                self._touch(kb_id)
                return self._loaded[kb_id]
            while len(self._loaded) >= self._max_loaded:
                evict_id = self._access_order.pop(0)
                del self._loaded[evict_id]
            data = self._load_from_disk(kb_id)
            self._loaded[kb_id] = data
            self._access_order.append(kb_id)
            return data

    def _load_from_disk(self, kb_id: str) -> dict:
        global embeddings
        kb_path = self.get_kb_path(kb_id)
        if not os.path.isdir(kb_path):
            raise FileNotFoundError(f"知识库 '{kb_id}' 不存在: {kb_path}")
        faiss_path = os.path.join(kb_path, "index.faiss")
        if not os.path.exists(faiss_path):
            raise FileNotFoundError(f"知识库 '{kb_id}' 尚未构建索引，请先上传文档并构建")
        vectorstore = FAISS.load_local(kb_path, embeddings, allow_dangerous_deserialization=True)
        retriever = vectorstore.as_retriever(search_kwargs={"k": HYBRID_CANDIDATE_K})
        bm25_index_obj = None
        bm25_docs_list = None
        bm25_path = os.path.join(kb_path, "bm25_index.pkl")
        if os.path.exists(bm25_path):
            with open(bm25_path, "rb") as f:
                bm25_data = pickle.load(f)
            bm25_index_obj = BM25Okapi(bm25_data["tokenized_corpus"])
            bm25_docs_list = bm25_data["documents"]
        return {
            "vectorstore": vectorstore,
            "bm25_index": bm25_index_obj,
            "bm25_docs": bm25_docs_list,
            "retriever": retriever,
        }

    def invalidate(self, kb_id: str):
        with self._lock:
            if kb_id in self._loaded:
                del self._loaded[kb_id]
                if kb_id in self._access_order:
                    self._access_order.remove(kb_id)

    def put(self, kb_id: str, data: dict):
        """将数据写入缓存，超出上限时自动淘汰最久未用的条目"""
        with self._lock:
            if kb_id in self._loaded:
                self._loaded[kb_id] = data
                self._touch(kb_id)
                return
            while len(self._loaded) >= self._max_loaded:
                evict_id = self._access_order.pop(0)
                del self._loaded[evict_id]
            self._loaded[kb_id] = data
            self._access_order.append(kb_id)

    def _touch(self, kb_id: str):
        if kb_id in self._access_order:
            self._access_order.remove(kb_id)
        self._access_order.append(kb_id)


def init_default_kb():
    """确保默认知识库目录结构和 registry.json 存在"""
    kb_root = os.path.join(base_dir, "knowledge_bases")
    default_kb = os.path.join(kb_root, "default")
    default_materials = os.path.join(default_kb, "materials")
    os.makedirs(default_materials, exist_ok=True)

    registry_path = os.path.join(kb_root, "registry.json")
    if not os.path.exists(registry_path):
        now = datetime.now().isoformat()
        registry = {
            "default": {
                "id": "default",
                "name": "默认知识库",
                "created_at": now,
                "updated_at": now,
                "description": ""
            }
        }
        with open(registry_path, "w", encoding="utf-8") as f:
            json.dump(registry, f, ensure_ascii=False, indent=2)
        print("创建知识库注册表: registry.json")


# 全局知识库管理器实例
kb_manager = None

# ============================================================
# SQLite 数据库函数
# ============================================================

def get_db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS sessions (
            session_id TEXT PRIMARY KEY,
            kb_id TEXT DEFAULT 'default',
            title TEXT DEFAULT '新对话',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            summary TEXT DEFAULT NULL
        );
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY (session_id) REFERENCES sessions(session_id) ON DELETE CASCADE
        );
        CREATE INDEX IF NOT EXISTS idx_messages_session_id ON messages(session_id);
    """)
    # 迁移：为已有数据库添加 kb_id 列
    columns = [row[1] for row in conn.execute("PRAGMA table_info(sessions)").fetchall()]
    if "kb_id" not in columns:
        conn.execute("ALTER TABLE sessions ADD COLUMN kb_id TEXT DEFAULT 'default'")
        print("数据库迁移：添加 sessions.kb_id 列")
    conn.commit()
    conn.close()
    print("数据库初始化完成")

def create_session(session_id: str, title: str = "新对话", kb_id: str = "default") -> dict:
    now = datetime.now().isoformat()
    conn = get_db()
    conn.execute(
        "INSERT INTO sessions (session_id, kb_id, title, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
        (session_id, kb_id, title, now, now)
    )
    conn.commit()
    conn.close()
    return {"session_id": session_id, "kb_id": kb_id, "title": title, "created_at": now, "updated_at": now}

def get_session(session_id: str) -> Optional[dict]:
    conn = get_db()
    row = conn.execute("SELECT * FROM sessions WHERE session_id = ?", (session_id,)).fetchone()
    conn.close()
    if row:
        return dict(row)
    return None

def list_sessions(kb_id: Optional[str] = None) -> List[dict]:
    conn = get_db()
    if kb_id:
        rows = conn.execute("""
            SELECT s.*, COUNT(m.id) as message_count
            FROM sessions s
            LEFT JOIN messages m ON s.session_id = m.session_id
            WHERE s.kb_id = ?
            GROUP BY s.session_id
            ORDER BY s.updated_at DESC
        """, (kb_id,)).fetchall()
    else:
        rows = conn.execute("""
            SELECT s.*, COUNT(m.id) as message_count
            FROM sessions s
            LEFT JOIN messages m ON s.session_id = m.session_id
            GROUP BY s.session_id
            ORDER BY s.updated_at DESC
        """).fetchall()
    conn.close()
    return [dict(r) for r in rows]

def delete_session(session_id: str) -> bool:
    conn = get_db()
    cursor = conn.execute("DELETE FROM sessions WHERE session_id = ?", (session_id,))
    deleted = cursor.rowcount > 0
    conn.commit()
    conn.close()
    return deleted

def add_message(session_id: str, role: str, content: str):
    now = datetime.now().isoformat()
    conn = get_db()
    conn.execute(
        "INSERT INTO messages (session_id, role, content, created_at) VALUES (?, ?, ?, ?)",
        (session_id, role, content, now)
    )
    conn.execute(
        "UPDATE sessions SET updated_at = ? WHERE session_id = ?",
        (now, session_id)
    )
    conn.commit()
    conn.close()

def get_recent_messages(session_id: str, limit: int = 20) -> List[dict]:
    conn = get_db()
    rows = conn.execute(
        "SELECT role, content, created_at FROM messages WHERE session_id = ? ORDER BY id DESC LIMIT ?",
        (session_id, limit)
    ).fetchall()
    conn.close()
    return [dict(r) for r in reversed(rows)]

def get_message_count(session_id: str) -> int:
    conn = get_db()
    row = conn.execute("SELECT COUNT(*) as cnt FROM messages WHERE session_id = ?", (session_id,)).fetchone()
    conn.close()
    return row["cnt"] if row else 0

def update_session_title(session_id: str, title: str):
    conn = get_db()
    conn.execute("UPDATE sessions SET title = ? WHERE session_id = ?", (title, session_id))
    conn.commit()
    conn.close()

def update_session_summary(session_id: str, summary: str):
    conn = get_db()
    conn.execute("UPDATE sessions SET summary = ? WHERE session_id = ?", (summary, session_id))
    conn.commit()
    conn.close()

def cleanup_old_sessions(max_age_days: int = SESSION_MAX_AGE_DAYS) -> int:
    cutoff = (datetime.now() - timedelta(days=max_age_days)).isoformat()
    conn = get_db()
    cursor = conn.execute("DELETE FROM sessions WHERE updated_at < ?", (cutoff,))
    deleted = cursor.rowcount
    conn.commit()
    conn.close()
    return deleted

# ============================================================
# 文档加载与分割
# ============================================================

def _load_pptx(file_path):
    """加载 .pptx 文件，每个幻灯片生成一个 Document。"""
    from pptx import Presentation
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            page_content = "\n".join(texts)
            docs.append(Document(
                page_content=page_content,
                metadata={"source": file_path, "page": slide_num}
            ))
    return docs


def _load_docx(file_path):
    """加载 .docx 文件，整个文档生成一个 Document（由下游文本分割器分块）。"""
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    page_content = "\n\n".join(paragraphs)
    return [Document(
        page_content=page_content,
        metadata={"source": file_path, "page": 1}
    )]


def load_documents(directory_path):
    """加载目录中的文档，支持 PDF、PPTX、DOCX、MD 格式。"""
    LOADER_MAP = {
        ".pdf": lambda fp: PyPDFLoader(fp).load(),
        ".pptx": _load_pptx,
        ".docx": _load_docx,
        ".md": lambda fp: TextLoader(fp, autodetect_encoding=True).load(),
    }

    all_documents = []
    stats = {}

    if not os.path.isdir(directory_path):
        print(f"错误：目录不存在 {directory_path}")
        return []

    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()
            if ext not in LOADER_MAP:
                continue

            file_path = os.path.join(root, filename)
            if ext not in stats:
                stats[ext] = {"loaded": 0, "skipped": 0, "errors": []}

            try:
                docs = LOADER_MAP[ext](file_path)
                if docs:
                    all_documents.extend(docs)
                    stats[ext]["loaded"] += 1
                    print(f"  [OK] {ext.upper()}: {filename} -> {len(docs)} 个文档")
                else:
                    stats[ext]["skipped"] += 1
                    print(f"  [SKIP] {ext.upper()}: {filename} (无有效文本)")
            except Exception as e:
                stats[ext]["skipped"] += 1
                stats[ext]["errors"].append(f"{filename}: {e}")
                print(f"  [ERROR] {ext.upper()}: {filename} -> {e}")

    print(f"\n文档加载汇总:")
    print(f"  总计加载: {len(all_documents)} 个文档片段")
    for ext, s in sorted(stats.items()):
        print(f"  {ext.upper()}: 加载={s['loaded']}, 跳过={s['skipped']}")
        for err in s["errors"]:
            print(f"    错误: {err}")

    return all_documents

def split_documents(documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=["\n\n", "\n", "。", "！", "？", "；", " ", ""]
    )
    texts = text_splitter.split_documents(documents)
    print(f"分割为 {len(texts)} 个文本块")
    return texts

# ============================================================
# 向量库创建与初始化（含 BM25 索引构建）
# ============================================================

def create_vectorstore(texts, kb_id="default"):
    global embedding_api_key

    if not embedding_api_key:
        print("错误：未设置EMBEDDING_API_KEY环境变量")
        return None

    from langchain_core.documents import Document

    if texts and not isinstance(texts[0], Document):
        print("警告：texts不是Document对象列表，尝试转换")
        new_texts = []
        for i, text in enumerate(texts):
            if isinstance(text, str):
                cleaned_text = str(text).strip()
                cleaned_text = ''.join(c for c in cleaned_text if ord(c) >= 32 or c in '\n\t')
                if cleaned_text:
                    new_texts.append(Document(page_content=cleaned_text, metadata={}))
                else:
                    print(f"跳过空文本元素: {i}")
            else:
                print(f"跳过非字符串元素: {type(text)}")
        texts = new_texts
    else:
        cleaned_texts = []
        for doc in texts:
            if hasattr(doc, 'page_content'):
                cleaned_content = str(doc.page_content).strip()
                cleaned_content = ''.join(c for c in cleaned_content if ord(c) >= 32 or c in '\n\t')
                if cleaned_content:
                    cleaned_texts.append(Document(page_content=cleaned_content, metadata=doc.metadata))
                else:
                    print("跳过空文本Document")
            else:
                print(f"跳过非Document对象: {type(doc)}")
        texts = cleaned_texts

    if not texts:
        print("错误：没有有效文本可处理")
        return None

    print(f"清理后文本块数量: {len(texts)}")

    emb = OpenAICompatibleEmbeddings()
    batch_size = 32
    vs = None

    for i in range(0, len(texts), batch_size):
        batch_texts = texts[i:i + batch_size]
        print(f"处理第 {i//batch_size + 1} 批文本块，共 {len(batch_texts)} 个")

        valid_docs = []
        for j, doc in enumerate(batch_texts):
            if hasattr(doc, 'page_content') and isinstance(doc.page_content, str) and doc.page_content.strip():
                valid_docs.append(doc)
            else:
                print(f"跳过无效的Document对象: {j}")

        if valid_docs:
            if vs is None:
                vs = FAISS.from_documents(valid_docs, emb)
            else:
                temp_vs = FAISS.from_documents(valid_docs, emb)
                vs.merge_from(temp_vs)

    if vs:
        kb_path = kb_manager.get_kb_path(kb_id)
        os.makedirs(kb_path, exist_ok=True)
        vs.save_local(kb_path)
        print(f"向量库已保存到 {kb_path}")

        # 构建 BM25 索引并持久化
        print("正在构建 BM25 索引...")
        tokenized_corpus = [tokenize(doc.page_content) for doc in texts]
        bm25_data = {"tokenized_corpus": tokenized_corpus, "documents": texts}
        bm25_path = os.path.join(kb_path, "bm25_index.pkl")
        with open(bm25_path, "wb") as f:
            pickle.dump(bm25_data, f)
        print(f"BM25 索引已保存到 {bm25_path}")

        return vs
    else:
        print("错误：向量库创建失败")
        return None


def _embed_texts(texts):
    """对清理后的 Document 列表做 embedding，返回 FAISS 索引（不保存到磁盘）。"""
    if not texts:
        return None
    emb = OpenAICompatibleEmbeddings()
    batch_size = 32
    vs = None
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        valid = [d for d in batch if hasattr(d, 'page_content') and d.page_content.strip()]
        if valid:
            if vs is None:
                vs = FAISS.from_documents(valid, emb)
            else:
                vs.merge_from(FAISS.from_documents(valid, emb))
    return vs


def _clean_docs(texts):
    """清洗 Document 列表，去除空内容和非法字符。"""
    from langchain_core.documents import Document
    cleaned = []
    for doc in texts:
        if hasattr(doc, 'page_content'):
            content = str(doc.page_content).strip()
            content = ''.join(c for c in content if ord(c) >= 32 or c in '\n\t')
            if content:
                cleaned.append(Document(page_content=content, metadata=doc.metadata))
    return cleaned


def _load_single_file(file_path):
    """加载单个文档文件，返回 Document 列表。"""
    LOADER_MAP = {
        ".pdf": lambda fp: PyPDFLoader(fp).load(),
        ".pptx": _load_pptx,
        ".docx": _load_docx,
        ".md": lambda fp: TextLoader(fp, autodetect_encoding=True).load(),
    }
    ext = os.path.splitext(file_path)[1].lower()
    loader = LOADER_MAP.get(ext)
    if not loader:
        return []
    try:
        docs = loader(file_path)
        return docs if docs else []
    except Exception as e:
        print(f"  [ERROR] 加载 {file_path}: {e}")
        return []


def load_manifest(kb_id="default"):
    """加载 manifest.json，返回 {filename: {mtime, chunks}} 字典。"""
    manifest_path = os.path.join(kb_manager.get_kb_path(kb_id), "manifest.json")
    if os.path.exists(manifest_path):
        with open(manifest_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {"files": {}}


def save_manifest(manifest, kb_id="default"):
    """保存 manifest.json。"""
    manifest_path = os.path.join(kb_manager.get_kb_path(kb_id), "manifest.json")
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)


def index_new_files(kb_id="default"):
    """增量索引：只处理指定知识库中新增或修改的文件，合并到现有索引。"""
    kb_path = kb_manager.get_kb_path(kb_id)
    materials_path = kb_manager.get_materials_path(kb_id)

    if not os.path.isdir(materials_path):
        return {"status": "success", "indexed": [], "message": "文档目录不存在"}

    if not os.path.isdir(kb_path):
        return {"status": "error", "message": "知识库不存在，请先调用 /init 初始化"}

    manifest = load_manifest(kb_id)
    indexed_files = manifest.get("files", {})

    # 有索引但无 manifest（旧数据兼容），需要全量重建
    if not indexed_files:
        print("增量索引：检测到旧索引无 manifest，执行全量重建...")
        success = init_vectorstore(kb_id=kb_id, force_rebuild=True)
        if success:
            return {"status": "success", "indexed": ["full_rebuild"], "message": "旧索引无 manifest，已全量重建"}
        return {"status": "error", "message": "全量重建失败"}

    # 扫描材料目录，找出新增或修改的文件
    SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".md"}
    new_files = []
    for filename in os.listdir(materials_path):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED_EXTS:
            continue
        fp = os.path.join(materials_path, filename)
        mtime = os.stat(fp).st_mtime
        if filename not in indexed_files or indexed_files[filename].get("mtime") != mtime:
            new_files.append(filename)

    if not new_files:
        return {"status": "success", "indexed": [], "message": "无新增或修改的文件"}

    print(f"增量索引：发现 {len(new_files)} 个新文件: {new_files}")

    # 加载知识库数据
    kb_data = kb_manager.get(kb_id)
    vs = kb_data["vectorstore"]
    bm25_docs_list = kb_data["bm25_docs"] or []

    indexed = []
    all_new_docs = []
    for filename in new_files:
        fp = os.path.join(materials_path, filename)
        mtime = os.stat(fp).st_mtime
        docs = _load_single_file(fp)
        if not docs:
            print(f"  [SKIP] {filename} (无有效文本)")
            continue
        texts = split_documents(docs)
        texts = _clean_docs(texts)
        if not texts:
            continue
        temp_vs = _embed_texts(texts)
        if temp_vs:
            vs.merge_from(temp_vs)
            all_new_docs.extend(texts)
            indexed.append({"file": filename, "chunks": len(texts)})
            indexed_files[filename] = {"mtime": mtime, "chunks": len(texts)}
            print(f"  [OK] {filename} -> {len(texts)} 个文本块")

    if indexed:
        vs.save_local(kb_path)
        print(f"向量库已保存，新增 {len(all_new_docs)} 个文本块")

        # 追加 BM25 索引
        bm25_docs_list.extend(all_new_docs)
        tokenized_corpus = [tokenize(d.page_content) for d in bm25_docs_list]
        bm25_path = os.path.join(kb_path, "bm25_index.pkl")
        with open(bm25_path, "wb") as f:
            pickle.dump({"tokenized_corpus": tokenized_corpus, "documents": bm25_docs_list}, f)

        manifest["files"] = indexed_files
        save_manifest(manifest, kb_id)
        kb_manager.invalidate(kb_id)

    return {"status": "success", "indexed": indexed, "message": f"增量索引完成，处理 {len(indexed)} 个文件"}


def remove_file_from_index(filename, kb_id="default"):
    """从 FAISS 和 BM25 索引中移除指定文件的所有向量。"""
    kb_path = kb_manager.get_kb_path(kb_id)

    if not os.path.isdir(kb_path):
        return {"status": "error", "message": "知识库不存在"}

    kb_data = kb_manager.get(kb_id)
    vs = kb_data["vectorstore"]
    bm25_docs_list = kb_data["bm25_docs"]
    bm25_idx = kb_data["bm25_index"]

    # 从 docstore 中查找属于该文件的所有 docstore ID
    ids_to_delete = []
    for idx, doc_id in vs.index_to_docstore_id.items():
        doc = vs.docstore.search(doc_id)
        if doc and os.path.basename(doc.metadata.get("source", "")) == filename:
            ids_to_delete.append(doc_id)

    if not ids_to_delete:
        return {"status": "success", "removed": 0, "message": f"索引中未找到 {filename} 的向量"}

    vs.delete(ids_to_delete)
    vs.save_local(kb_path)
    print(f"已从 FAISS 索引中移除 {filename} 的 {len(ids_to_delete)} 个向量")

    # 重建 BM25 索引（排除被删除文件的文档）
    if bm25_docs_list is not None:
        bm25_docs_list = [d for d in bm25_docs_list if os.path.basename(d.metadata.get("source", "")) != filename]
        tokenized_corpus = [tokenize(d.page_content) for d in bm25_docs_list]
        bm25_path = os.path.join(kb_path, "bm25_index.pkl")
        with open(bm25_path, "wb") as f:
            pickle.dump({"tokenized_corpus": tokenized_corpus, "documents": bm25_docs_list}, f)
        print(f"BM25 索引已更新（剩余 {len(bm25_docs_list)} 个文档）")

    # 更新 manifest
    manifest = load_manifest(kb_id)
    if filename in manifest.get("files", {}):
        del manifest["files"][filename]
        save_manifest(manifest, kb_id)

    kb_manager.invalidate(kb_id)
    return {"status": "success", "removed": len(ids_to_delete), "message": f"已移除 {filename} 的 {len(ids_to_delete)} 个向量"}


def init_shared_models():
    """初始化共享的 embeddings、LLM、reranker（仅需执行一次）"""
    global embeddings, llm, reranker
    global embedding_api_key, llm_api_key, llm_base_url, llm_model
    try:
        if not embedding_api_key:
            print("未设置EMBEDDING_API_KEY环境变量")
            return False

        if not llm_api_key:
            print("未设置LLM_API_KEY环境变量")
            return False

        from langchain_openai import ChatOpenAI

        embeddings = OpenAICompatibleEmbeddings()
        llm = ChatOpenAI(
            openai_api_key=llm_api_key,
            openai_api_base=llm_base_url,
            model_name=llm_model,
            temperature=0.3,
            request_timeout=120,
            max_retries=2
        )
        reranker = OpenAICompatibleReranker()
        print(f"重排序器已初始化 (模型: {RERANKER_MODEL})")
        return True
    except Exception as e:
        print(f"初始化共享模型失败: {e}")
        return False


def init_vectorstore(kb_id="default", force_rebuild=False):
    """初始化指定知识库的向量库和 BM25 索引"""
    try:
        if not embeddings:
            print("共享模型未初始化，请先调用 init_shared_models()")
            return False

        kb_path = kb_manager.get_kb_path(kb_id)
        materials_path = kb_manager.get_materials_path(kb_id)

        if force_rebuild and os.path.exists(kb_path):
            # 只删除索引文件，保留 materials 目录
            for fname in ["index.faiss", "index.pkl", "bm25_index.pkl", "manifest.json"]:
                fpath = os.path.join(kb_path, fname)
                if os.path.exists(fpath):
                    os.remove(fpath)
            print("已删除旧的索引文件，准备重新创建")

        if not force_rebuild and os.path.exists(kb_path):
            vectorstore = FAISS.load_local(kb_path, embeddings, allow_dangerous_deserialization=True)
            retriever = vectorstore.as_retriever(search_kwargs={"k": HYBRID_CANDIDATE_K})
            print("成功加载现有向量库")

            # 加载 BM25 索引
            bm25_index_obj = None
            bm25_docs_list = None
            bm25_path = os.path.join(kb_path, "bm25_index.pkl")
            if os.path.exists(bm25_path):
                with open(bm25_path, "rb") as f:
                    bm25_data = pickle.load(f)
                bm25_index_obj = BM25Okapi(bm25_data["tokenized_corpus"])
                bm25_docs_list = bm25_data["documents"]
                print(f"BM25 索引已加载 ({len(bm25_docs_list)} 个文档)")
            else:
                print("警告：BM25 索引文件不存在，混合检索不可用。请使用 /init?force_rebuild=true 重建知识库")

            # 存入缓存
            kb_manager.put(kb_id, {
                "vectorstore": vectorstore,
                "bm25_index": bm25_index_obj,
                "bm25_docs": bm25_docs_list,
                "retriever": retriever,
            })

            # 加载 manifest
            manifest = load_manifest(kb_id)
            if manifest.get("files"):
                print(f"manifest 已加载 ({len(manifest['files'])} 个文件)")
            else:
                print("manifest 为空，将在下次增量索引时全量重建")
        else:
            try:
                docs = load_documents(materials_path)
                if docs:
                    texts = split_documents(docs)
                    vs = create_vectorstore(texts, kb_id)

                    # 构建 manifest
                    chunk_counts = {}
                    for doc in texts:
                        fname = os.path.basename(doc.metadata.get("source", ""))
                        if fname:
                            chunk_counts[fname] = chunk_counts.get(fname, 0) + 1
                    manifest = {"files": {}}
                    for fname, count in chunk_counts.items():
                        fp = os.path.join(materials_path, fname)
                        if os.path.isfile(fp):
                            manifest["files"][fname] = {"mtime": os.stat(fp).st_mtime, "chunks": count}
                    save_manifest(manifest, kb_id)
                    print(f"manifest 已构建 ({len(manifest['files'])} 个文件)")
                    print("从文档创建了新的向量库和 BM25 索引")
                else:
                    print("未找到文档，向量库未初始化")
                    return False
            except Exception as e:
                print(f"创建向量库失败: {e}")
                return False

            # 将新建的索引直接写入缓存，避免从磁盘重新加载
            if vs:
                retriever = vs.as_retriever(search_kwargs={"k": HYBRID_CANDIDATE_K})
                # 从 create_vectorstore 已保存的 pickle 加载 BM25，保证数据一致性
                bm25_index_obj = None
                bm25_docs_list = None
                bm25_path = os.path.join(kb_path, "bm25_index.pkl")
                if os.path.exists(bm25_path):
                    with open(bm25_path, "rb") as f:
                        bm25_data = pickle.load(f)
                    bm25_index_obj = BM25Okapi(bm25_data["tokenized_corpus"])
                    bm25_docs_list = bm25_data["documents"]
                kb_manager.put(kb_id, {
                    "vectorstore": vs,
                    "bm25_index": bm25_index_obj,
                    "bm25_docs": bm25_docs_list,
                    "retriever": retriever,
                })

        return True
    except Exception as e:
        print(f"初始化向量库失败: {e}")
        return False

# ============================================================
# 预检索优化：查询改写
# ============================================================

def rewrite_query(question: str, session_id: Optional[str] = None) -> str:
    """用 LLM 将口语化问题改写为规范化检索词"""
    history_text = format_history(session_id) if session_id else ""
    if not history_text:
        history_text = "（无对话历史）\n"

    prompt = QUERY_REWRITE_PROMPT.format(history=history_text, question=question)
    try:
        answer = llm.invoke(prompt)
        rewritten = answer.content if hasattr(answer, 'content') else str(answer)
        rewritten = rewritten.strip().strip('"').strip("'")
        print(f"查询改写: '{question}' -> '{rewritten}'")
        return rewritten if rewritten else question
    except Exception as e:
        print(f"查询改写失败: {e}，使用原始问题")
        return question

# ============================================================
# 预检索优化：HyDE
# ============================================================

def hyde_generate(question: str) -> str:
    """让 LLM 生成假设性答案，用于检索"""
    prompt = HYDE_PROMPT.format(question=question)
    try:
        answer = llm.invoke(prompt)
        hypothetical = answer.content if hasattr(answer, 'content') else str(answer)
        print(f"HyDE 假设答案生成成功 (长度: {len(hypothetical)})")
        return hypothetical
    except Exception as e:
        print(f"HyDE 生成失败: {e}，使用原始问题")
        return question


def mqe_generate(question: str, n: int = 3) -> List[str]:
    """让 LLM 生成 n 个查询变体用于多查询扩展"""
    prompt = MQE_PROMPT.format(question=question, n=n)
    try:
        answer = llm.invoke(prompt)
        raw = answer.content if hasattr(answer, 'content') else str(answer)
        queries = [line.strip() for line in raw.strip().split('\n') if line.strip()]
        seen = set()
        result = []
        for q in queries:
            if q and q not in seen:
                seen.add(q)
                result.append(q)
        if not result:
            result = [question]
        print(f"MQE 生成 {len(result)} 个查询变体: {result}")
        return result
    except Exception as e:
        print(f"MQE 生成失败: {e}，使用原始问题")
        return [question]


# ============================================================
# 混合检索策略（向量 + BM25 + RRF 融合）
# ============================================================

def hybrid_retrieve(query: str, k: int = DEFAULT_TOP_K, kb_id: str = "default") -> list:
    """混合检索：FAISS 向量检索 + BM25 关键词检索，倒数排名融合"""
    kb_data = kb_manager.get(kb_id)
    vs = kb_data["vectorstore"]
    bm25_idx = kb_data["bm25_index"]
    bm25_docs_list = kb_data["bm25_docs"]
    candidate_k = HYBRID_CANDIDATE_K

    # 1. 向量检索
    vector_docs = vs.similarity_search(query, k=candidate_k)

    # 2. BM25 检索
    bm25_results = []
    if bm25_idx is not None and bm25_docs_list is not None:
        query_tokens = tokenize(query)
        scores = bm25_idx.get_scores(query_tokens)
        top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:candidate_k]
        bm25_results = [bm25_docs_list[i] for i in top_indices if scores[i] > 0]

    # 3. 倒数排名融合（Reciprocal Rank Fusion）
    doc_scores = {}
    rrf_k = 60  # RRF 常数

    for rank, doc in enumerate(vector_docs):
        key = doc.page_content[:200]  # 用前200字符作为去重key
        doc_scores[key] = doc_scores.get(key, {"doc": doc, "score": 0})
        doc_scores[key]["score"] += VECTOR_WEIGHT / (rank + rrf_k)

    for rank, doc in enumerate(bm25_results):
        key = doc.page_content[:200]
        doc_scores[key] = doc_scores.get(key, {"doc": doc, "score": 0})
        doc_scores[key]["score"] += BM25_WEIGHT / (rank + rrf_k)

    # 按融合分数排序
    sorted_items = sorted(doc_scores.values(), key=lambda x: x["score"], reverse=True)
    result = []
    for item in sorted_items[:k]:
        doc = item["doc"]
        if hasattr(doc, 'metadata'):
            doc.metadata['rrf_score'] = round(item["score"], 6)
        result.append(doc)
    print(f"混合检索: 向量={len(vector_docs)}篇, BM25={len(bm25_results)}篇, 融合后={len(result)}篇")
    return result

# ============================================================
# RAG 查询函数（带检索优化 + 会话历史）
# ============================================================

def format_history(session_id: str) -> str:
    if not session_id:
        return ""
    messages = get_recent_messages(session_id, limit=MAX_HISTORY_EXCHANGES * 2)
    if not messages:
        return ""
    lines = []
    for msg in messages:
        role_label = "用户" if msg["role"] == "user" else "助手"
        lines.append(f"{role_label}：{msg['content']}")
    return "### 对话历史：\n" + "\n\n".join(lines) + "\n\n"

def extract_sources(docs, top_k=DEFAULT_TOP_K):
    """从检索文档中提取来源信息"""
    sources = []
    for i, doc in enumerate(docs[:top_k]):
        meta = doc.metadata if hasattr(doc, 'metadata') else {}
        source = {
            "rank": i + 1,
            "content": doc.page_content[:500] if hasattr(doc, 'page_content') else str(doc)[:500],
            "source": meta.get("source", ""),
            "page": meta.get("page", None),
        }
        if "vector_score" in meta:
            source["vector_score"] = meta["vector_score"]
        if "rrf_score" in meta:
            source["rrf_score"] = meta["rrf_score"]
        if "rerank_score" in meta:
            source["rerank_score"] = meta["rerank_score"]
        sources.append(source)
    return sources


def rag_query(question: str, session_id: Optional[str] = None,
              kb_id: str = "default",
              retrieval_strategy: str = "default",
              pre_retrieval=None,
              post_retrieval: str = "none",
              top_k: Optional[int] = None) -> tuple:
    prompt_text, sources = _retrieve_and_build_prompt(
        question, session_id, kb_id, retrieval_strategy,
        pre_retrieval, post_retrieval, use_history=True, top_k=top_k
    )
    answer = llm.invoke(prompt_text)
    answer_text = answer.content if hasattr(answer, 'content') else str(answer)
    return answer_text, sources

def rag_query_stateless(question: str,
                        kb_id: str = "default",
                        retrieval_strategy: str = "default",
                        pre_retrieval=None,
                        post_retrieval: str = "none",
                        top_k: Optional[int] = None) -> tuple:
    prompt_text, sources = _retrieve_and_build_prompt(
        question, None, kb_id, retrieval_strategy,
        pre_retrieval, post_retrieval, use_history=False, top_k=top_k
    )
    answer = llm.invoke(prompt_text)
    answer_text = answer.content if hasattr(answer, 'content') else str(answer)
    return answer_text, sources

def _retrieve_and_build_prompt(question, session_id=None, kb_id="default",
                               retrieval_strategy="default",
                               pre_retrieval=None, post_retrieval="none",
                               use_history=True, top_k=None, history_text=None):
    """公共检索逻辑，返回 (prompt_text, sources)"""
    if pre_retrieval is None:
        pre_retrieval = []
    effective_top_k = max(1, min(top_k, 20)) if top_k is not None else DEFAULT_TOP_K

    kb_data = kb_manager.get(kb_id)
    vs = kb_data["vectorstore"]

    # --- Phase 1: 顺序执行预检索查询变换 (rewrite -> hyde -> mqe) ---
    current_query = question

    if "rewrite" in pre_retrieval:
        current_query = rewrite_query(current_query, session_id) if session_id else rewrite_query(current_query)

    if "hyde" in pre_retrieval:
        current_query = hyde_generate(current_query)

    mqe_queries = None
    if "mqe" in pre_retrieval:
        mqe_queries = mqe_generate(current_query)
        if current_query not in mqe_queries:
            mqe_queries.insert(0, current_query)

    # --- Phase 2: 检索 ---
    if mqe_queries is not None:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        def _retrieve_one(q):
            if "hyde" in pre_retrieval and retrieval_strategy == "hybrid":
                # HyDE + Hybrid: 用 HyDE 向量做向量检索 + BM25 关键词检索融合
                q_emb = embeddings.embed_query(q)
                vector_results = vs.similarity_search_with_score_by_vector(q_emb, k=HYBRID_CANDIDATE_K)
                vector_docs = [doc for doc, _ in vector_results]
                bm25_results = []
                kb_data_inner = kb_manager.get(kb_id)
                bm25_idx_inner = kb_data_inner["bm25_index"]
                bm25_docs_inner = kb_data_inner["bm25_docs"]
                if bm25_idx_inner is not None and bm25_docs_inner is not None:
                    q_tokens = tokenize(q)
                    scores = bm25_idx_inner.get_scores(q_tokens)
                    top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:HYBRID_CANDIDATE_K]
                    bm25_results = [bm25_docs_inner[i] for i in top_indices if scores[i] > 0]
                # RRF 融合
                doc_scores = {}
                rrf_k_local = 60
                for r, doc in enumerate(vector_docs):
                    key = doc.page_content[:200]
                    doc_scores[key] = doc_scores.get(key, {"doc": doc, "score": 0})
                    doc_scores[key]["score"] += VECTOR_WEIGHT / (r + rrf_k_local)
                for r, doc in enumerate(bm25_results):
                    key = doc.page_content[:200]
                    doc_scores[key] = doc_scores.get(key, {"doc": doc, "score": 0})
                    doc_scores[key]["score"] += BM25_WEIGHT / (r + rrf_k_local)
                sorted_items = sorted(doc_scores.values(), key=lambda x: x["score"], reverse=True)
                return [item["doc"] for item in sorted_items[:HYBRID_CANDIDATE_K]]
            elif "hyde" in pre_retrieval:
                q_emb = embeddings.embed_query(q)
                scored = vs.similarity_search_with_score_by_vector(q_emb, k=HYBRID_CANDIDATE_K)
                return [doc for doc, _ in scored]
            elif retrieval_strategy == "hybrid":
                return hybrid_retrieve(q, k=HYBRID_CANDIDATE_K, kb_id=kb_id)
            else:
                scored = vs.similarity_search_with_score(q, k=HYBRID_CANDIDATE_K)
                return [doc for doc, _ in scored]

        all_docs = {}
        rrf_k = 60
        with ThreadPoolExecutor(max_workers=min(len(mqe_queries), 4)) as executor:
            futures = {executor.submit(_retrieve_one, q): i for i, q in enumerate(mqe_queries)}
            for future in as_completed(futures):
                q_idx = futures[future]
                q_docs = future.result()
                for rank, doc in enumerate(q_docs):
                    key = doc.page_content[:200]
                    if key not in all_docs:
                        all_docs[key] = {"doc": doc, "score": 0}
                    all_docs[key]["score"] += 1.0 / (rank + rrf_k)

        sorted_items = sorted(all_docs.values(), key=lambda x: x["score"], reverse=True)
        docs = []
        for item in sorted_items[:HYBRID_CANDIDATE_K]:
            doc = item["doc"]
            if hasattr(doc, 'metadata'):
                doc.metadata['rrf_score'] = round(item["score"], 6)
            docs.append(doc)

    elif "hyde" in pre_retrieval:
        hyde_embedding = embeddings.embed_query(current_query)
        scored = vs.similarity_search_with_score_by_vector(hyde_embedding, k=HYBRID_CANDIDATE_K)
        docs = []
        for doc, score in scored:
            if hasattr(doc, 'metadata'):
                doc.metadata['vector_score'] = round(float(score), 4)
            docs.append(doc)

    elif retrieval_strategy == "hybrid":
        docs = hybrid_retrieve(current_query, k=HYBRID_CANDIDATE_K, kb_id=kb_id)

    else:
        scored = vs.similarity_search_with_score(current_query, k=HYBRID_CANDIDATE_K)
        docs = []
        for doc, score in scored:
            if hasattr(doc, 'metadata'):
                doc.metadata['vector_score'] = round(float(score), 4)
            docs.append(doc)

    # --- Phase 3: 后检索 ---
    if post_retrieval == "rerank" and reranker:
        docs = reranker.rerank(current_query, docs, top_n=effective_top_k)

    # --- Phase 4: 截断并构建 prompt ---
    top_docs = docs[:effective_top_k]
    context = "\n\n".join([doc.page_content for doc in top_docs])
    sources = extract_sources(docs, top_k=effective_top_k)

    if history_text:
        prompt_text = PROMPT_WITH_HISTORY.format(context=context, history=history_text, question=question)
    elif use_history and session_id:
        history_text = format_history(session_id)
        prompt_text = PROMPT_WITH_HISTORY.format(context=context, history=history_text, question=question)
    else:
        prompt_text = PROMPT.format(context=context, question=question)

    return prompt_text, sources

def rag_query_stream(question, session_id=None, kb_id="default",
                     retrieval_strategy="default",
                     pre_retrieval=None, post_retrieval="none", top_k=None):
    """流式 RAG 查询（同步生成器）：先 yield sources JSON，再逐 token yield"""
    prompt_text, sources = _retrieve_and_build_prompt(
        question, session_id, kb_id, retrieval_strategy,
        pre_retrieval, post_retrieval, use_history=True, top_k=top_k
    )
    yield json.dumps({"type": "sources", "data": sources}, ensure_ascii=False)
    for chunk in llm.stream(prompt_text):
        content = chunk.content if hasattr(chunk, 'content') else str(chunk)
        if content:
            yield json.dumps({"type": "token", "data": content}, ensure_ascii=False)

def rag_query_stateless_stream(question, kb_id="default",
                               retrieval_strategy="default",
                               pre_retrieval=None, post_retrieval="none", top_k=None):
    """无状态流式 RAG 查询（同步生成器）"""
    prompt_text, sources = _retrieve_and_build_prompt(
        question, None, kb_id, retrieval_strategy,
        pre_retrieval, post_retrieval, use_history=False, top_k=top_k
    )
    yield json.dumps({"type": "sources", "data": sources}, ensure_ascii=False)
    for chunk in llm.stream(prompt_text):
        content = chunk.content if hasattr(chunk, 'content') else str(chunk)
        if content:
            yield json.dumps({"type": "token", "data": content}, ensure_ascii=False)

def generate_summary_task(session_id: str):
    """后台任务：生成会话摘要"""
    try:
        messages = get_recent_messages(session_id, limit=999)
        if len(messages) < 2:
            return
        conversation = "\n".join(
            f"{'用户' if m['role'] == 'user' else '助手'}：{m['content']}"
            for m in messages
        )
        prompt = SUMMARY_PROMPT_TEMPLATE.format(conversation=conversation)
        answer = llm.invoke(prompt)
        summary = answer.content if hasattr(answer, 'content') else str(answer)
        update_session_summary(session_id, summary)
        print(f"会话 {session_id} 摘要已生成")
    except Exception as e:
        print(f"生成摘要失败: {e}")

# ============================================================
# API 端点
# ============================================================

@app.get("/")
async def root():
    return {"message": "知识库API已启动（检索优化版），支持查询改写/HyDE/混合检索/重排序"}

@app.get("/health")
async def health_check():
    return {"status": "healthy"}


def mask_key(key):
    """Mask API key, showing only last 4 characters"""
    if not key or len(key) < 8:
        return "***" if key else ""
    return "***" + key[-4:]


def save_env_file():
    """Save current config to .env file, preserving comments and formatting"""
    env_path = os.path.join(base_dir, ".env")
    updates = {
        "SILICONFLOW_API_KEY": siliconflow_api_key or "",
        "LLM_API_KEY": llm_api_key or "",
        "LLM_BASE_URL": llm_base_url,
        "LLM_MODEL": llm_model,
        "EMBEDDING_API_KEY": embedding_api_key or "",
        "EMBEDDING_BASE_URL": embedding_base_url,
        "EMBEDDING_MODEL": embedding_model_name,
        "RERANKER_MODEL": RERANKER_MODEL,
        "RERANKER_API_KEY": reranker_api_key or "",
        "RERANKER_BASE_URL": reranker_base_url,
    }
    updated_keys = set()

    lines = []
    if os.path.exists(env_path):
        with open(env_path, "r", encoding="utf-8") as f:
            for line in f:
                stripped = line.strip()
                if stripped and not stripped.startswith("#") and "=" in stripped:
                    key = stripped.split("=", 1)[0].strip()
                    if key in updates:
                        lines.append(f"{key}={updates[key]}\n")
                        updated_keys.add(key)
                        continue
                lines.append(line)

    # Append any vars not found in the original file
    for key, val in updates.items():
        if key not in updated_keys:
            lines.append(f"{key}={val}\n")

    with open(env_path, "w", encoding="utf-8") as f:
        f.writelines(lines)


@app.get("/config")
async def get_config():
    return {
        "siliconflow_api_key": mask_key(siliconflow_api_key) if siliconflow_api_key else "",
        "llm": {
            "model": llm_model,
            "base_url": llm_base_url,
            "api_key": mask_key(llm_api_key),
        },
        "embedding": {
            "model": embedding_model_name,
            "base_url": embedding_base_url,
            "api_key": mask_key(embedding_api_key),
        },
        "reranker": {
            "model": RERANKER_MODEL,
            "base_url": reranker_base_url,
            "api_key": mask_key(reranker_api_key),
        },
    }


@app.post("/config")
async def update_config(request: dict):
    global llm_model, llm_base_url, llm_api_key
    global embedding_base_url, embedding_api_key, embedding_model_name
    global llm, embeddings, reranker, RERANKER_MODEL
    global reranker_base_url, reranker_api_key

    global siliconflow_api_key

    sf_key = request.get("siliconflow_api_key", "")
    if sf_key and not sf_key.startswith("***"):
        siliconflow_api_key = sf_key

    # 优先使用新格式 models，回退到顶层 llm/embedding/reranker
    models = request.get("models", {})
    llm_cfg = models.get("llm") or request.get("llm", {})
    emb_cfg = models.get("embedding") or request.get("embedding", {})
    rer_cfg = models.get("reranker") or request.get("reranker", {})

    def resolve_key(cfg, existing_key):
        """解析 API key：显式指定 > 硅基流动共享 key > 保留现有"""
        k = cfg.get("api_key", "")
        if k and not k.startswith("***"):
            return k
        # 硅基流动 URL 自动用共享 key
        base = cfg.get("base_url") or ""
        if "siliconflow" in base and siliconflow_api_key:
            return siliconflow_api_key
        return existing_key

    def resolve_base_url(cfg, existing):
        return cfg.get("base_url") or existing

    # 更新 LLM
    if llm_cfg.get("model"): llm_model = llm_cfg["model"]
    llm_base_url = resolve_base_url(llm_cfg, llm_base_url)
    llm_api_key = resolve_key(llm_cfg, llm_api_key)

    # 更新 Embedding
    if emb_cfg.get("model"): embedding_model_name = emb_cfg["model"]
    embedding_base_url = resolve_base_url(emb_cfg, embedding_base_url)
    embedding_api_key = resolve_key(emb_cfg, embedding_api_key)

    # 更新 Reranker
    if rer_cfg.get("model"): RERANKER_MODEL = rer_cfg["model"]
    reranker_base_url = resolve_base_url(rer_cfg, reranker_base_url)
    reranker_api_key = resolve_key(rer_cfg, reranker_api_key)

    print(f"[配置更新] LLM: {llm_model} @ {llm_base_url}")
    print(f"[配置更新] Embedding: {embedding_model_name} @ {embedding_base_url}")
    print(f"[配置更新] Reranker: {RERANKER_MODEL} @ {reranker_base_url}")

    # Reinitialize components
    errors = []
    try:
        from langchain_openai import ChatOpenAI
        llm = ChatOpenAI(
            openai_api_key=llm_api_key,
            openai_api_base=llm_base_url,
            model_name=llm_model,
            temperature=0.3,
            request_timeout=120,
            max_retries=2,
        )
    except Exception as e:
        errors.append(f"LLM 初始化失败: {e}")

    try:
        embeddings = OpenAICompatibleEmbeddings()
    except Exception as e:
        errors.append(f"Embedding 初始化失败: {e}")

    try:
        reranker = OpenAICompatibleReranker(model=RERANKER_MODEL)
    except Exception as e:
        errors.append(f"Reranker 初始化失败: {e}")

    # Persist to .env
    try:
        save_env_file()
    except Exception as e:
        errors.append(f"保存 .env 失败: {e}")

    if errors:
        return {"status": "partial", "message": "; ".join(errors)}

    return {"status": "success", "message": "配置已更新并保存"}


# SiliconFlow 免费模型列表（硬编码兜底，同时支持从 API 动态获取）
SILICONFLOW_KNOWN_FREE_MODELS = {
    "llm": [
        "Qwen/Qwen3-8B",
        "Qwen/Qwen3-14B",
        "Qwen/Qwen3-30B-A3B-Instruct-2507",
        "Qwen/Qwen3.5-4B",
        "Qwen/Qwen3.5-9B",
        "Qwen/Qwen3.5-27B",
        "Qwen/Qwen3.5-35B-A3B",
        "Qwen/Qwen2.5-7B-Instruct",
        "THUDM/GLM-4-9B-0414",
        "THUDM/GLM-Z1-9B-0414",
        "deepseek-ai/DeepSeek-V3",
        "deepseek-ai/DeepSeek-R1",
        "deepseek-ai/DeepSeek-V4-Flash",
        "zai-org/GLM-4.5-Air",
    ],
    "embedding": [
        "BAAI/bge-large-zh-v1.5",
        "BAAI/bge-large-en-v1.5",
        "BAAI/bge-m3",
        "netease-youdao/bce-embedding-base_v1",
    ],
    "reranker": [
        "BAAI/bge-reranker-v2-m3",
        "netease-youdao/bce-reranker-base_v1",
    ],
}


@app.get("/models/siliconflow")
async def get_siliconflow_models(x_api_key: Optional[str] = Header(None, alias="X-API-Key")):
    """获取硅基流动可用模型列表，按类型分类。"""
    key = x_api_key or embedding_api_key
    if not key:
        return {"status": "success", "models": SILICONFLOW_KNOWN_FREE_MODELS, "source": "hardcoded"}

    try:
        resp = requests.get(
            "https://api.siliconflow.cn/v1/models",
            headers={"Authorization": f"Bearer {key}"},
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        all_models = [m["id"] for m in data.get("data", [])]

        llm, emb, rerank = [], [], []
        for mid in all_models:
            low = mid.lower()
            # 跳过 Pro/ 前缀（付费）和 LoRA/ 前缀
            if mid.startswith("Pro/") or mid.startswith("LoRA/"):
                continue
            if "rerank" in low:
                rerank.append(mid)
            elif "embed" in low or (low.startswith("bge-") and "rerank" not in low):
                emb.append(mid)
            elif any(kw in low for kw in [
                "instruct", "chat", "-v3", "-v4", "-r1", "qwen", "glm",
                "llama", "internlm", "minimax", "yi-", "baichuan",
                "deepseek", "seed-oss", "hunyuan-a13b", "step-",
                "ling-flash", "ling-mini", "kimi",
            ]):
                llm.append(mid)

        models = {
            "llm": sorted(llm),
            "embedding": sorted(emb),
            "reranker": sorted(rerank),
        }
        return {"status": "success", "models": models, "source": "api"}
    except Exception as e:
        print(f"获取硅基流动模型列表失败: {e}")
        return {"status": "success", "models": SILICONFLOW_KNOWN_FREE_MODELS, "source": "hardcoded"}


@app.post("/init")
async def init_knowledge_base(force_rebuild: bool = False, kb_id: str = "default"):
    try:
        global embedding_api_key
        if not embedding_api_key:
            return {"status": "error", "message": "未设置EMBEDDING_API_KEY环境变量"}

        success = await asyncio.to_thread(init_vectorstore, kb_id=kb_id, force_rebuild=force_rebuild)
        if success:
            message = f"知识库 '{kb_id}' 重建成功" if force_rebuild else f"知识库 '{kb_id}' 初始化成功"
            return {"status": "success", "message": message}
        else:
            return {"status": "error", "message": f"知识库 '{kb_id}' 初始化失败，请确保材料文件夹中有支持的文件"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".md"}

@app.get("/materials")
async def list_materials(kb_id: str = "default"):
    """列出指定知识库材料目录中的文档文件。"""
    try:
        materials_path = kb_manager.get_materials_path(kb_id)
        if not os.path.isdir(materials_path):
            return {"status": "success", "files": []}

        files = []
        for name in os.listdir(materials_path):
            ext = os.path.splitext(name)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                continue
            fp = os.path.join(materials_path, name)
            stat = os.stat(fp)
            files.append({
                "name": name,
                "size": stat.st_size,
                "modified": stat.st_mtime,
            })
        files.sort(key=lambda f: f["modified"], reverse=True)
        return {"status": "success", "files": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/materials/upload")
async def upload_material(files: List[UploadFile] = File(...), kb_id: str = "default"):
    """上传文档到指定知识库的材料目录。"""
    try:
        materials_path = kb_manager.get_materials_path(kb_id)
        os.makedirs(materials_path, exist_ok=True)

        uploaded = []
        errors = []
        for upload_file in files:
            safe_name = os.path.basename(upload_file.filename)
            if not safe_name:
                errors.append(f"{upload_file.filename}: 无效文件名")
                continue
            ext = os.path.splitext(safe_name)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                errors.append(f"{safe_name}: 不支持的文件格式（仅支持 PDF/PPTX/DOCX/MD）")
                continue
            dest = os.path.join(materials_path, safe_name)
            content = await upload_file.read()
            with open(dest, "wb") as f:
                f.write(content)
            uploaded.append(safe_name)

        if errors:
            return {"status": "partial" if uploaded else "error", "uploaded": uploaded, "errors": errors}

        # 上传成功后自动增量索引
        index_result = await asyncio.to_thread(index_new_files, kb_id=kb_id)
        return {"status": "success", "uploaded": uploaded, "index": index_result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/materials/index")
async def index_materials(kb_id: str = "default"):
    """手动触发增量索引，只处理新增或修改的文件。"""
    try:
        result = await asyncio.to_thread(index_new_files, kb_id=kb_id)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/materials/{filename}")
async def delete_material(filename: str, kb_id: str = "default"):
    """删除指定知识库材料目录中的指定文件。"""
    try:
        filename = os.path.normpath(filename)
        if os.path.dirname(filename) != '':
            raise HTTPException(status_code=400, detail="非法文件名")

        materials_path = kb_manager.get_materials_path(kb_id)
        fp = os.path.join(materials_path, filename)
        if not os.path.realpath(fp).startswith(os.path.realpath(materials_path) + os.sep):
            raise HTTPException(status_code=400, detail="非法文件名")
        if not os.path.isfile(fp):
            raise HTTPException(status_code=404, detail=f"文件不存在: {filename}")

        os.remove(fp)

        # 从索引中移除该文件的向量
        index_result = await asyncio.to_thread(remove_file_from_index, filename, kb_id)

        return {"status": "success", "message": f"已删除: {filename}", "index": index_result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/ask/stream")
async def ask_stream(query: Query, background_tasks: BackgroundTasks):
    """流式问答接口，返回 SSE 事件流"""
    global embedding_api_key, llm_api_key
    if not llm_api_key:
        raise HTTPException(status_code=503, detail="未设置LLM_API_KEY环境变量")
    if not embedding_api_key:
        raise HTTPException(status_code=503, detail="未设置EMBEDDING_API_KEY环境变量")

    # 解析 kb_id：优先用请求参数，其次从会话获取
    kb_id = query.kb_id
    session_id = query.session_id
    if session_id:
        session = get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="会话不存在")
        if not kb_id:
            kb_id = session.get("kb_id", "default")
    if not kb_id:
        kb_id = "default"

    def event_generator():
        full_answer = []
        try:
            gen = rag_query_stream if session_id else rag_query_stateless_stream
            kwargs = {
                "question": query.question,
                "kb_id": kb_id,
                "retrieval_strategy": query.retrieval_strategy,
                "pre_retrieval": query.pre_retrieval,
                "post_retrieval": query.post_retrieval,
                "top_k": query.top_k,
            }
            if session_id:
                kwargs["session_id"] = session_id

            for event_str in gen(**kwargs):
                event = json.loads(event_str)
                if event["type"] == "token":
                    full_answer.append(event["data"])
                yield f"data: {event_str}\n\n"

            yield "data: [DONE]\n\n"
        except Exception as e:
            print(f"[ask/stream 错误] {e}")
            yield f"data: {json.dumps({'type': 'error', 'data': str(e)}, ensure_ascii=False)}\n\n"
        finally:
            # 流式完成后保存消息（finally 保证客户端断开时也能保存）
            answer_text = "".join(full_answer)
            if session_id and answer_text:
                try:
                    add_message(session_id, "user", query.question)
                    add_message(session_id, "assistant", answer_text)
                    msg_count = get_message_count(session_id)
                    if msg_count == 2:
                        title = query.question[:30] + ("..." if len(query.question) > 30 else "")
                        update_session_title(session_id, title)
                    if msg_count > 0 and msg_count % SUMMARY_TRIGGER_COUNT == 0:
                        background_tasks.add_task(generate_summary_task, session_id)
                except Exception as save_err:
                    print(f"[ask/stream 保存消息失败] {save_err}")

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/ask", response_model=Response)
async def ask_question(query: Query, background_tasks: BackgroundTasks):
    try:
        global embedding_api_key, llm_api_key
        if not llm_api_key:
            raise HTTPException(status_code=503, detail="未设置LLM_API_KEY环境变量")
        if not embedding_api_key:
            raise HTTPException(status_code=503, detail="未设置EMBEDDING_API_KEY环境变量")

        # 解析 kb_id
        kb_id = query.kb_id
        session_id = query.session_id
        if session_id:
            session = get_session(session_id)
            if not session:
                raise HTTPException(status_code=404, detail="会话不存在")
            if not kb_id:
                kb_id = session.get("kb_id", "default")
        if not kb_id:
            kb_id = "default"

        strategy_info = f"retrieval={query.retrieval_strategy}, pre={query.pre_retrieval}, post={query.post_retrieval}"
        print(f"收到问题: {query.question} (kb={kb_id}, session={session_id}, {strategy_info})")

        try:
            if session_id:
                answer, sources = await asyncio.wait_for(
                    asyncio.to_thread(
                        rag_query, query.question, session_id, kb_id,
                        query.retrieval_strategy, query.pre_retrieval, query.post_retrieval,
                        query.top_k
                    ),
                    timeout=180
                )
            else:
                answer, sources = await asyncio.wait_for(
                    asyncio.to_thread(
                        rag_query_stateless, query.question, kb_id,
                        query.retrieval_strategy, query.pre_retrieval, query.post_retrieval,
                        query.top_k
                    ),
                    timeout=180
                )
        except asyncio.TimeoutError:
            print("LLM调用超时(180秒)")
            raise HTTPException(status_code=504, detail="LLM调用超时，请检查模型配置或网络连接")

        print(f"LLM返回结果，长度: {len(answer)}")

        if session_id:
            add_message(session_id, "user", query.question)
            add_message(session_id, "assistant", answer)

            msg_count = get_message_count(session_id)
            if msg_count == 2:
                title = query.question[:30] + ("..." if len(query.question) > 30 else "")
                update_session_title(session_id, title)

            if msg_count > 0 and msg_count % SUMMARY_TRIGGER_COUNT == 0:
                background_tasks.add_task(generate_summary_task, session_id)

        return Response(answer=answer, sources=sources, session_id=session_id)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/v1/models")
async def list_models():
    return OpenAIModelsResponse(
        data=[OpenAIModel(id="knowledge-base")]
    )

@app.post("/v1/chat/completions")
async def chat_completions(request: OpenAIChatRequest, background_tasks: BackgroundTasks):
    global embedding_api_key, llm_api_key
    if not llm_api_key:
        raise HTTPException(status_code=503, detail="未设置LLM_API_KEY环境变量")
    if not embedding_api_key:
        raise HTTPException(status_code=503, detail="未设置EMBEDDING_API_KEY环境变量")

    # 解析 kb_id
    kb_id = request.kb_id
    session_id = request.session_id
    if session_id:
        session = get_session(session_id)
        if session and not kb_id:
            kb_id = session.get("kb_id", "default")
    if not kb_id:
        kb_id = "default"

    user_message = ""
    for msg in reversed(request.messages):
        if msg.role == "user":
            user_message = msg.content
            break

    if not user_message:
        raise HTTPException(status_code=400, detail="No user message found")

    if session_id:
        session = get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="会话不存在")

    print(f"[chat] 收到问题: {user_message} (session_id={session_id}, stream={request.stream})")

    # 流式输出
    if request.stream:
        import time as _time
        chat_id = f"chatcmpl-{uuid.uuid4().hex[:12]}"

        def openai_stream_generator():
            full_answer = []
            try:
                gen = rag_query_stream if session_id else rag_query_stateless_stream
                kwargs = {
                    "question": user_message,
                    "kb_id": kb_id,
                    "retrieval_strategy": request.retrieval_strategy,
                    "pre_retrieval": request.pre_retrieval,
                    "post_retrieval": request.post_retrieval,
                    "top_k": request.top_k,
                }
                if session_id:
                    kwargs["session_id"] = session_id

                for event_str in gen(**kwargs):
                    event = json.loads(event_str)
                    if event["type"] == "token":
                        full_answer.append(event["data"])
                        chunk = {
                            "id": chat_id,
                            "object": "chat.completion.chunk",
                            "created": int(_time.time()),
                            "model": request.model,
                            "choices": [{"index": 0, "delta": {"content": event["data"]}, "finish_reason": None}],
                        }
                        yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"

                # 结束 chunk
                chunk = {
                    "id": chat_id,
                    "object": "chat.completion.chunk",
                    "created": int(_time.time()),
                    "model": request.model,
                    "choices": [{"index": 0, "delta": {}, "finish_reason": "stop"}],
                }
                yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
                yield "data: [DONE]\n\n"
            except Exception as e:
                error_chunk = {"error": {"message": str(e), "type": "server_error"}}
                yield f"data: {json.dumps(error_chunk, ensure_ascii=False)}\n\n"
            finally:
                # 保存消息（finally 保证客户端断开时也能保存）
                answer_text = "".join(full_answer)
                if session_id and answer_text:
                    try:
                        add_message(session_id, "user", user_message)
                        add_message(session_id, "assistant", answer_text)
                        msg_count = get_message_count(session_id)
                        if msg_count == 2:
                            title = user_message[:30] + ("..." if len(user_message) > 30 else "")
                            update_session_title(session_id, title)
                        if msg_count > 0 and msg_count % SUMMARY_TRIGGER_COUNT == 0:
                            background_tasks.add_task(generate_summary_task, session_id)
                    except Exception as save_err:
                        print(f"[chat/stream 保存消息失败] {save_err}")

        return StreamingResponse(
            openai_stream_generator(),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    # 非流式输出
    try:
        if session_id:
            answer, _sources = await asyncio.wait_for(
                asyncio.to_thread(
                    rag_query, user_message, session_id, kb_id,
                    request.retrieval_strategy, request.pre_retrieval, request.post_retrieval,
                    request.top_k
                ),
                timeout=180
            )
        else:
            if len(request.messages) > 1:
                history_lines = []
                for msg in request.messages[:-1]:
                    role_label = "用户" if msg.role == "user" else "助手"
                    history_lines.append(f"{role_label}：{msg.content}")
                history_text = "### 对话历史：\n" + "\n\n".join(history_lines) + "\n\n"

                prompt_text, _sources = _retrieve_and_build_prompt(
                    user_message, None, kb_id,
                    request.retrieval_strategy, request.pre_retrieval, request.post_retrieval,
                    use_history=False, top_k=request.top_k, history_text=history_text
                )
                answer_obj = await asyncio.wait_for(
                    asyncio.to_thread(llm.invoke, prompt_text),
                    timeout=180
                )
                answer = answer_obj.content if hasattr(answer_obj, 'content') else str(answer_obj)
            else:
                answer, _sources = await asyncio.wait_for(
                    asyncio.to_thread(
                        rag_query_stateless, user_message, kb_id,
                        request.retrieval_strategy, request.pre_retrieval, request.post_retrieval,
                        request.top_k
                    ),
                    timeout=180
                )
    except asyncio.TimeoutError:
        print("[chat] LLM调用超时(180秒)")
        raise HTTPException(status_code=504, detail="LLM调用超时，请检查模型配置或网络连接")

    if session_id:
        add_message(session_id, "user", user_message)
        add_message(session_id, "assistant", answer)
        msg_count = get_message_count(session_id)
        if msg_count == 2:
            title = user_message[:30] + ("..." if len(user_message) > 30 else "")
            update_session_title(session_id, title)
        if msg_count > 0 and msg_count % SUMMARY_TRIGGER_COUNT == 0:
            background_tasks.add_task(generate_summary_task, session_id)

    return OpenAIChatResponse(
        model=request.model,
        choices=[
            OpenAIChatChoice(
                message=OpenAIMessage(role="assistant", content=answer)
            )
        ]
    )

# ============================================================
# 会话管理端点
# ============================================================

@app.post("/sessions", response_model=SessionInfo)
async def create_new_session(body: SessionCreate = None):
    title = body.title if body and body.title else "新对话"
    kb_id = body.kb_id if body and body.kb_id else "default"
    session_id = str(uuid.uuid4())
    session = create_session(session_id, title, kb_id)
    return SessionInfo(
        session_id=session["session_id"],
        kb_id=session["kb_id"],
        title=session["title"],
        created_at=session["created_at"],
        updated_at=session["updated_at"],
        message_count=0
    )

@app.get("/sessions", response_model=SessionListResponse)
async def get_all_sessions(kb_id: Optional[str] = None):
    sessions = list_sessions(kb_id=kb_id)
    return SessionListResponse(
        sessions=[
            SessionInfo(
                session_id=s["session_id"],
                kb_id=s.get("kb_id", "default"),
                title=s["title"],
                created_at=s["created_at"],
                updated_at=s["updated_at"],
                summary=s.get("summary"),
                message_count=s.get("message_count", 0)
            )
            for s in sessions
        ]
    )

@app.get("/sessions/{session_id}", response_model=SessionInfo)
async def get_session_info(session_id: str):
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="会话不存在")
    msg_count = get_message_count(session_id)
    return SessionInfo(
        session_id=session["session_id"],
        kb_id=session.get("kb_id", "default"),
        title=session["title"],
        created_at=session["created_at"],
        updated_at=session["updated_at"],
        summary=session.get("summary"),
        message_count=msg_count
    )

@app.delete("/sessions/{session_id}")
async def delete_session_endpoint(session_id: str):
    deleted = delete_session(session_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="会话不存在")
    return {"status": "deleted", "session_id": session_id}

@app.get("/sessions/{session_id}/messages", response_model=SessionHistoryResponse)
async def get_session_messages(session_id: str):
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="会话不存在")
    messages = get_recent_messages(session_id, limit=999)
    return SessionHistoryResponse(
        session_id=session_id,
        messages=[
            MessageInfo(role=m["role"], content=m["content"], created_at=m["created_at"])
            for m in messages
        ]
    )

@app.post("/sessions/cleanup")
async def cleanup_sessions_endpoint(max_age_days: int = SESSION_MAX_AGE_DAYS):
    deleted = cleanup_old_sessions(max_age_days)
    return {"deleted_sessions": deleted, "max_age_days": max_age_days}

# ============================================================
# 知识库管理 API
# ============================================================

@app.get("/knowledge-bases")
async def list_knowledge_bases():
    return {"knowledge_bases": kb_manager.list_kbs()}

@app.post("/knowledge-bases")
async def create_knowledge_base(kb: KBCreate):
    info = kb_manager.create_kb(kb.name, kb.description)
    return info

@app.get("/knowledge-bases/{kb_id}")
async def get_knowledge_base(kb_id: str):
    registry = kb_manager.load_registry()
    if kb_id not in registry:
        raise HTTPException(status_code=404, detail=f"知识库 '{kb_id}' 不存在")
    kbs = kb_manager.list_kbs()
    for kb in kbs:
        if kb["id"] == kb_id:
            return kb
    return registry[kb_id]

@app.put("/knowledge-bases/{kb_id}")
async def update_knowledge_base(kb_id: str, kb: KBCreate):
    if not kb_manager.update_kb(kb_id, name=kb.name, description=kb.description):
        raise HTTPException(status_code=404, detail=f"知识库 '{kb_id}' 不存在")
    return {"message": "更新成功"}

@app.delete("/knowledge-bases/{kb_id}")
async def delete_knowledge_base(kb_id: str):
    if kb_id == "default":
        raise HTTPException(status_code=400, detail="不能删除默认知识库")
    if not kb_manager.delete_kb(kb_id):
        raise HTTPException(status_code=404, detail=f"知识库 '{kb_id}' 不存在")
    return {"message": f"知识库 '{kb_id}' 已删除"}

@app.post("/knowledge-bases/{kb_id}/rebuild")
async def rebuild_knowledge_base(kb_id: str):
    result = await asyncio.to_thread(init_vectorstore, kb_id=kb_id, force_rebuild=True)
    if result:
        return {"message": f"知识库 '{kb_id}' 重建成功"}
    raise HTTPException(status_code=500, detail="知识库重建失败")

# ============================================================
# 启动
# ============================================================

init_default_kb()
kb_manager = KBManager(max_loaded=2)
init_db()
init_shared_models()

# 清理 knowledge_bases/ 根目录下的孤立 index 文件（FAISS 旧 bug 产生的乱码文件）
_kb_root = os.path.join(base_dir, "knowledge_bases")
if os.path.isdir(_kb_root):
    for _f in os.listdir(_kb_root):
        if _f.endswith('.faiss') or _f.endswith('.faiss.pkl'):
            try:
                os.remove(os.path.join(_kb_root, _f))
                print(f"已清理孤立索引文件: {_f}")
            except Exception:
                pass

if __name__ == "__main__":
    print("=" * 50)
    print("启动知识库API（检索优化版）...")
    print(f"硅基流动 API Key: {'已设置' if siliconflow_api_key else '未设置'}")
    print(f"LLM: {llm_model} @ {llm_base_url}")
    print(f"Embedding: {embedding_model_name} @ {embedding_base_url}")
    print(f"Reranker: {RERANKER_MODEL} @ {reranker_base_url}")
    print(f"数据库路径: {DB_PATH}")
    print(f"知识库数量: {len(kb_manager.list_kbs())}")
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=8001, log_level="info")
